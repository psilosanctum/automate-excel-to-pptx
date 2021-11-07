import pandas as pd
import numpy as np
import datetime as dt
from pptx import Presentation
from pptx.util import Inches, Pt


data_path = 'weather_data.xlsx'

template_path = Presentation(
    'template_weather_report.pptx')


def createTable(dataframe, slideTitle):
    dataframe = dataframe.reset_index()
    dataframe['timestamp'] = dataframe['timestamp'].dt.strftime(
        '%m/%d/%y')

    arrayOfDataframes = []

    for i, (index, row) in enumerate(dataframe.iterrows()):
        if (row.name % 6 == 0) & (row.name != 0):
            cutDataFrame = dataframe[index - 6: row.name]
            cutDataFrame = cutDataFrame.reset_index(drop=True)
            arrayOfDataframes.append(cutDataFrame)
        if (row.name % 6 == 0) & ((row.name == len(dataframe) - 1) | (row.name == len(dataframe) - 2) | (row.name == len(dataframe) - 3) | (row.name == len(dataframe) - 4) | (row.name == len(dataframe) - 5) | (row.name == len(dataframe) - 6)):
            cutDataFrame = dataframe[index: len(dataframe)]
            cutDataFrame = cutDataFrame.reset_index(drop=True)
            arrayOfDataframes.append(cutDataFrame)
        if (len(dataframe) - 1 == 2) | (len(dataframe) - 2 == 1) | (len(dataframe) - 3 == 2) | (len(dataframe) - 4 == 1) | (len(dataframe) - 5 == 2):
            cutDataFrame = dataframe[0: len(dataframe)]
            cutDataFrame = cutDataFrame.reset_index(drop=True)
            arrayOfDataframes = []
            arrayOfDataframes.append(cutDataFrame)

    for i in arrayOfDataframes:
        title_only_slide_layout = template_path.slide_layouts[1]
        slide = template_path.slides.add_slide(title_only_slide_layout)
        shapes = slide.shapes

        shapes.title.text = slideTitle
        rows = i.shape[0] + 1
        cols = len(dataframe.columns)
        left = Inches(0.0)
        top = Inches(1.0)
        width = Inches(7.5)
        height = Inches(3.5)

        table = shapes.add_table(
            rows, cols, left, top, width, height).table

        # Set column widths
        table.columns[0].width = Inches(1.50)
        table.columns[1].width = Inches(1.80)
        table.columns[2].width = Inches(1.25)
        table.columns[3].width = Inches(1.75)
        table.columns[4].width = Inches(2.00)
        table.columns[5].width = Inches(1.50)
        table.columns[6].width = Inches(2.2)
        table.columns[7].width = Inches(1.35)

        # Set column headings
        table.cell(0, 0).text = 'Conditions'
        table.cell(0, 1).text = 'Dew PointC'
        table.cell(0, 2).text = 'Humidity'
        table.cell(0, 3).text = 'TemperatureC'
        table.cell(0, 4).text = 'Wind Direction'
        table.cell(0, 5).text = 'Visibility Km'
        table.cell(0, 6).text = 'Wind Dir Degrees'
        table.cell(0, 7).text = 'timestamp'

        for index, row in i.iterrows():
            table.cell(index + 1, 0).text = row['Conditions']
            table.cell(index + 1, 1).text = row['Dew PointC']
            table.cell(index + 1, 2).text = row['Humidity']
            table.cell(index + 1, 3).text = row['TemperatureC']
            table.cell(index + 1, 4).text = row['Wind Direction']
            table.cell(index + 1, 5).text = row['VisibilityKm']
            table.cell(index + 1, 6).text = row['WindDirDegrees']
            table.cell(index + 1, 7).text = row['timestamp']

    def iter_cells(table):
        for row in table.rows:
            for cell in row.cells:
                yield cell

        for cell in iter_cells(table):
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(10)

    template_path.save(
        'automated_weather_report.pptx')


def onlyShowScatteredCloudWest():
    dataframe = pd.read_excel(open(data_path, 'rb'),
                              sheet_name='weather6')
    dataframe = dataframe[['Conditions', 'Dew PointC', 'Humidity',
                           'TemperatureC', 'Wind Direction', 'VisibilityKm', 'WindDirDegrees', 'timestamp']].astype(str)
    dataframe['timestamp'] = pd.to_datetime(dataframe['timestamp'])
    dataframe = dataframe.set_index(['timestamp'])
    dataframe = dataframe.iloc[:, [0, 1, 2, 3, 4, 5, 6]]
    dataframe = dataframe.loc['2015']
    dataframe = dataframe.fillna(value='Zero')
    dataframe = dataframe[(dataframe['Conditions'] == 'Scattered Clouds')]
    dataframe = dataframe[(dataframe['Dew PointC'] == '7')]
    dataframe = dataframe[(dataframe['Wind Direction'] == 'West')]
    dataframe = dataframe.sort_values(['Humidity'])
    createTable(dataframe, "Scattered Clouds, West Winds & \n Dew Point of 7")


def onlyShowCalmHaze():
    dataframe = pd.read_excel(open(data_path, 'rb'),
                              sheet_name='weather6')
    dataframe = dataframe[['Conditions', 'Dew PointC', 'Humidity',
                           'TemperatureC', 'Wind Direction', 'VisibilityKm', 'WindDirDegrees', 'timestamp']].astype(str)
    dataframe['timestamp'] = pd.to_datetime(dataframe['timestamp'])
    dataframe = dataframe.set_index(['timestamp'])
    dataframe = dataframe.iloc[:, [0, 1, 2, 3, 4, 5, 6]]
    dataframe = dataframe.loc['2015']
    dataframe = dataframe.fillna(value='Zero')
    dataframe = dataframe[(dataframe['Conditions'] == 'Haze')]
    dataframe = dataframe[(dataframe['Dew PointC'] == '5')]
    dataframe = dataframe[(dataframe['Wind Direction'] == 'Calm')]
    dataframe = dataframe.sort_values(['Humidity'])
    createTable(dataframe, "Haze, Calm Winds & \n Dew Point of 5")


def runAllTasks():
    onlyShowCalmHaze()
    onlyShowScatteredCloudWest()


runAllTasks()
